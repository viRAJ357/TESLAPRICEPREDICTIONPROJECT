import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # 0 is Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Tesla Stock Price Prediction"
subtitle.text = "Deep Learning for Time Series Forecasting\nProject Report for Labmentix Submit\nJanuary 2026"

# Slide 2: Introduction and Problem Statement
slide_layout = prs.slide_layouts[1] # 1 is Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction & Problem Statement"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Objective: Forecast future Tesla stock prices using historical data."
p = tf.add_paragraph()
p.text = "Challenges: Standard machine learning struggles with sequential temporal dependencies."
p = tf.add_paragraph()
p.text = "Approach: Deep learning designed for sequences - RNNs and LSTMs."
p = tf.add_paragraph()
p.text = "Target: Predict Adj Close price 1, 5, and 10 days into the future."

# Slide 3: Approach and Methodology - Data Preprocessing
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Preprocessing & EDA"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Feature Selection: 'Adj Close' as target; 'Date' as index."
p = tf.add_paragraph()
p.text = "Missing Values: Forward/backward-fill interpolation."
p = tf.add_paragraph()
p.text = "Data Scaling: MinMaxScaler (normalization between 0 and 1) for neural network stability."
p = tf.add_paragraph()
p.text = "Sequence Generation: Sliding window of past n days (e.g., 60) to predict day n+1."

# Slide 4: Deep Learning Modeling
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Deep Learning Architecture"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Model 1: SimpleRNN (Basic recurrent neural network)"
p = tf.add_paragraph()
p.text = "Model 2: LSTM (Advanced variant overcoming vanishing gradients, better long-term memory)"
p = tf.add_paragraph()
p.text = "Architecture specifics:"
p.level = 0
p1 = tf.add_paragraph()
p1.text = "Sequential architectures with Dropout layers to prevent overfitting."
p1.level = 1
p2 = tf.add_paragraph()
p2.text = "Dense layer for final single predicted price output."
p2.level = 1
p3 = tf.add_paragraph()
p3.text = "Optimizer: Adam | Loss Function: MSE | EarlyStopping utilized."
p3.level = 1

# Slide 5: Results and Evaluation
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Results & Evaluation"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Evaluation Metric: RMSE (Root Mean Squared Error) and MSE."
p = tf.add_paragraph()
p.text = "Comparison:"
p.level = 0
p1 = tf.add_paragraph()
p1.text = "LSTM outperformed SimpleRNN due to ability to retain long-term historical context."
p1.level = 1
p2 = tf.add_paragraph()
p2.text = "SimpleRNNs prone to forgetting earlier data points in long sequences."
p2.level = 1

# Slide 6: Insights and Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Insights & Conclusion"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Effectiveness: DL models (especially LSTM) capture macroeconomic trends and momentum well."
p = tf.add_paragraph()
p.text = "Limitations:"
p.level = 0
p1 = tf.add_paragraph()
p1.text = "Market fluctuations based on exogenous variables (news, earnings)."
p1.level = 1
p2 = tf.add_paragraph()
p2.text = "Lag Effect: Predictions can sometimes act as a slightly shifted version of previous day's price."
p2.level = 1

# Slide 7: Suggestions for Improvements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Suggestions for Improvements"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Alternative Data: Sentiment analysis (Twitter, Reddit, News)."
p = tf.add_paragraph()
p.text = "Multivariate Forecasting: Include Volume, Moving Averages, RSI, MACD."
p = tf.add_paragraph()
p.text = "Macro Indicators: Interest rates, inflation, S&P 500."
p = tf.add_paragraph()
p.text = "Advanced Architectures: GRU or Transformer models."

# Slide 8: Timeline and Deliverables
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Timeline & Deliverables"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Completed: Data Cleaning, EDA, Feature Engineering."
p = tf.add_paragraph()
p.text = "Completed: Deep Learning Modelling and Evaluation."
p = tf.add_paragraph()
p.text = "Final Submission Date: January 12, 2026."
p = tf.add_paragraph()
p.text = "Deliverables: Jupyter Notebook (execution code), Final Report, and Presentation."

# Save PPT
prs.save('Labmentix_Presentation.pptx')
print("Successfully generated Labmentix_Presentation.pptx")
